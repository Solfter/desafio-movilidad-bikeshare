"""
Reemplaza los marcadores FIG_1..FIG_5 de la presentación por las figuras Plotly
exportadas desde `01_eda_final.ipynb`.

USO
---
1) Exportar las figuras del notebook (ver `exportar_figuras.py`) a la carpeta `figuras/`.
2) Ejecutar:

       pip install python-pptx pillow
       python insertar_figuras.py

   Opcional, si tus archivos están en otra parte:

       python insertar_figuras.py --ppt presentacion_bikeshare.pptx \
                                  --figuras figuras/ \
                                  --salida presentacion_bikeshare_final.pptx

Cada figura se centra dentro del marco del marcador respetando su proporción:
no se deforma. Si un PNG no existe, ese marcador se deja intacto (así puedes
insertar las figuras de a poco, a medida que las tengas).
"""

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# Marcador  ->  archivo esperado dentro de la carpeta `figuras/`
MAPEO = {
    "FIG_1": "01_distribucion_demanda.png",   # Lámina 6  · Hernán
    "FIG_2": "02_heatmap_hora_dia.png",       # Lámina 7  · Hernán
    "FIG_3": "03_lluvia_y_feriados.png",      # Lámina 8  · Hernán
    "FIG_4": "04_autocorrelacion.png",        # Lámina 9  · Matías
    "FIG_5": "05_dashboard.png",              # Lámina 12 · Michelangelo (captura del dashboard)
}


def borrar(shape):
    """Elimina una forma del slide (python-pptx no expone un .delete())."""
    shape._element.getparent().remove(shape._element)


def encajar(marco, img_path):
    """Calcula posición y tamaño para meter la imagen dentro del marco sin deformarla."""
    with Image.open(img_path) as im:
        prop_img = im.width / im.height
    prop_marco = marco["w"] / marco["h"]

    if prop_img >= prop_marco:          # imagen más ancha: limita el ancho
        w = marco["w"]
        h = int(w / prop_img)
    else:                               # imagen más alta: limita el alto
        h = marco["h"]
        w = int(h * prop_img)

    x = marco["x"] + (marco["w"] - w) // 2
    y = marco["y"] + (marco["h"] - h) // 2
    return x, y, w, h


def main():
    ap = argparse.ArgumentParser(description="Inserta las figuras Plotly en la presentación.")
    ap.add_argument("--ppt", default="presentacion_bikeshare.pptx")
    ap.add_argument("--figuras", default="figuras")
    ap.add_argument("--salida", default="presentacion_bikeshare_final.pptx")
    args = ap.parse_args()

    carpeta = Path(args.figuras)
    prs = Presentation(args.ppt)
    insertadas, pendientes = [], []

    for n_slide, slide in enumerate(prs.slides, start=1):
        for marcador, archivo in MAPEO.items():
            # El marco del marcador: la forma que se llama exactamente como el marcador
            marco_shape = next((sh for sh in slide.shapes if sh.name == marcador), None)
            if marco_shape is None:
                continue

            ruta = carpeta / archivo
            if not ruta.exists():
                pendientes.append(f"{marcador} (lámina {n_slide}) → falta {ruta}")
                continue

            marco = {"x": marco_shape.left, "y": marco_shape.top,
                     "w": marco_shape.width, "h": marco_shape.height}

            # Borramos el marco y todas sus etiquetas (FIG_n, FIG_n_icono, FIG_n_pie, ...)
            for sh in [s for s in slide.shapes if s.name.startswith(marcador)]:
                borrar(sh)

            x, y, w, h = encajar(marco, ruta)
            slide.shapes.add_picture(str(ruta), Emu(x), Emu(y), Emu(w), Emu(h))
            insertadas.append(f"{marcador} (lámina {n_slide}) ← {archivo}")

    prs.save(args.salida)

    print(f"\nGuardado: {args.salida}\n")
    if insertadas:
        print("Figuras insertadas:")
        for i in insertadas:
            print(f"  ✔ {i}")
    if pendientes:
        print("\nMarcadores que quedaron sin figura (se dejaron intactos):")
        for p in pendientes:
            print(f"  · {p}")


if __name__ == "__main__":
    main()
